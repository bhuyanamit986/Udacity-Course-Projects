#!/usr/bin/env python3
"""
TED Talk-style PowerPoint presentation generator for Angular Optimization
Creates a visually engaging, cinematic presentation with dark theme and accent colors
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
import os

def create_angular_presentation():
    """Create the Angular Optimization TED Talk presentation"""
    
    # Create presentation
    prs = Presentation()
    
    # Set slide dimensions to widescreen (16:9)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Define color scheme based on template instructions
    colors = {
        # Template-specified colors
        'title_bg_start': RGBColor(26, 26, 46),    # #1a1a2e
        'title_bg_end': RGBColor(22, 33, 62),      # #16213e
        'content_bg': RGBColor(255, 255, 255),     # White background for content slides
        'angular_red': RGBColor(221, 0, 49),       # #DD0031 Angular Red
        'lighthouse_blue': RGBColor(66, 133, 244), # #4285F4 Lighthouse Blue
        'performance_green': RGBColor(52, 168, 83), # #34A853 Performance Green
        'light_blue': RGBColor(135, 206, 235),     # #87CEEB Light Blue
        'dark_blue': RGBColor(26, 26, 46),         # #1a1a2e Dark Blue
        'dark_gray': RGBColor(51, 51, 51),         # #333333 Dark Gray
        'white': RGBColor(255, 255, 255),
        'light_gray': RGBColor(200, 200, 200),
        
        # Slide-specific gradient colors
        'red_gradient_start': RGBColor(255, 107, 107),  # #FF6B6B
        'red_gradient_end': RGBColor(255, 142, 142),    # #FF8E8E
        'blue_gradient_start': RGBColor(66, 133, 244),  # #4285F4
        'blue_gradient_end': RGBColor(135, 206, 235),   # #87CEEB
        'green_gradient_start': RGBColor(52, 168, 83),  # #34A853
        'green_gradient_end': RGBColor(144, 238, 144),  # #90EE90
        'orange_gradient_start': RGBColor(251, 188, 4), # #FBBC04
        'orange_gradient_end': RGBColor(255, 215, 0),   # #FFD700
        'purple_gradient_start': RGBColor(156, 39, 176), # #9C27B0
        'purple_gradient_end': RGBColor(225, 190, 231),  # #E1BEE7
        'teal_gradient_start': RGBColor(0, 150, 136),    # #009688
        'teal_gradient_end': RGBColor(128, 203, 196),    # #80CBC4
        'indigo_gradient_start': RGBColor(63, 81, 181),  # #3F51B5
        'indigo_gradient_end': RGBColor(159, 168, 218),  # #9FA8DA
        'dark_gradient_start': RGBColor(44, 62, 80),     # #2C3E50
        'dark_gradient_end': RGBColor(52, 73, 94),       # #34495E
        'success_gradient_start': RGBColor(39, 174, 96), # #27AE60
        'success_gradient_end': RGBColor(88, 214, 141),  # #58D68D
        'future_gradient_start': RGBColor(142, 68, 173), # #8E44AD
        'future_gradient_end': RGBColor(187, 143, 206),  # #BB8FCE
        'action_gradient_start': RGBColor(231, 76, 60),  # #E74C3C
        'action_gradient_end': RGBColor(241, 148, 138),  # #F1948A
        'discussion_gradient_start': RGBColor(52, 152, 219), # #3498DB
        'discussion_gradient_end': RGBColor(133, 193, 233), # #85C1E9
        'gratitude_gradient_start': RGBColor(243, 156, 18), # #F39C12
        'gratitude_gradient_end': RGBColor(247, 220, 111)  # #F7DC6F
    }
    
    # Slide 1: Title Slide
    create_title_slide(prs, colors)
    
    # Slide 2: The 3-Second Rule Hook
    create_hook_slide(prs, colors)
    
    # Slide 3: Understanding Lighthouse Metrics
    create_lighthouse_metrics_slide(prs, colors)
    
    # Slide 4: Optimization Approach
    create_approach_slide(prs, colors)
    
    # Slide 5: Lazy Loading Deep Dive
    create_lazy_loading_slide(prs, colors)
    
    # Slide 6: OnPush Change Detection
    create_onpush_slide(prs, colors)
    
    # Slide 7: Font Loading Strategy
    create_font_loading_slide(prs, colors)
    
    # Slide 8: Tree-Shakeable Material Modules
    create_material_modules_slide(prs, colors)
    
    # Slide 9: Live Demo
    create_demo_slide(prs, colors)
    
    # Slide 10: Before/After Results
    create_results_slide(prs, colors)
    
    # Slide 11: Future Roadmap
    create_roadmap_slide(prs, colors)
    
    # Slide 12: Call to Action
    create_cta_slide(prs, colors)
    
    # Slide 13: Q&A
    create_qa_slide(prs, colors)
    
    # Slide 14: Thank You
    create_thank_you_slide(prs, colors)
    
    return prs

def create_title_slide(prs, colors):
    """Create the title slide with template-specified design"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set background to gradient from #1a1a2e to #16213e
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['title_bg_start']
    
    # Main title - Roboto Bold, 48px, White
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_p = title_frame.paragraphs[0]
    title_p.text = "The Zen of Angular Optimization: A Lighthouse-Guided Journey"
    title_p.font.size = Pt(48)
    title_p.font.bold = True
    title_p.font.name = 'Roboto'
    title_p.font.color.rgb = colors['white']
    title_p.alignment = PP_ALIGN.CENTER
    
    # Subtitle - Roboto Regular, 24px, Light Blue (#87CEEB)
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11.33), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.clear()
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.text = "From 4.5MB bundles to blazing-fast load times"
    subtitle_p.font.size = Pt(24)
    subtitle_p.font.name = 'Roboto'
    subtitle_p.font.color.rgb = colors['light_blue']
    subtitle_p.alignment = PP_ALIGN.CENTER
    
    # Add Angular and Lighthouse logo placeholders (top left and top right)
    # Angular logo placeholder (top left)
    angular_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2), Inches(0.8))
    angular_frame = angular_box.text_frame
    angular_frame.clear()
    angular_p = angular_frame.paragraphs[0]
    angular_p.text = "Angular"
    angular_p.font.size = Pt(20)
    angular_p.font.bold = True
    angular_p.font.color.rgb = colors['angular_red']
    
    # Lighthouse logo placeholder (top right)
    lighthouse_box = slide.shapes.add_textbox(Inches(10.83), Inches(0.5), Inches(2), Inches(0.8))
    lighthouse_frame = lighthouse_box.text_frame
    lighthouse_frame.clear()
    lighthouse_p = lighthouse_frame.paragraphs[0]
    lighthouse_p.text = "Lighthouse"
    lighthouse_p.font.size = Pt(20)
    lighthouse_p.font.bold = True
    lighthouse_p.font.color.rgb = colors['lighthouse_blue']
    lighthouse_p.alignment = PP_ALIGN.RIGHT
    
    # Performance indicator
    perf_box = slide.shapes.add_textbox(Inches(4), Inches(5.2), Inches(5.33), Inches(0.8))
    perf_frame = perf_box.text_frame
    perf_frame.clear()
    perf_p = perf_frame.paragraphs[0]
    perf_p.text = "Performance Score: 45 → 78"
    perf_p.font.size = Pt(20)
    perf_p.font.color.rgb = colors['performance_green']
    perf_p.alignment = PP_ALIGN.CENTER
    
    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = "Hook: Start with '3 seconds. That's all you have.' Make eye contact, pause for dramatic effect. Frame this as an opportunity, not a problem."

def create_hook_slide(prs, colors):
    """Create the 3-second rule hook slide with red gradient background"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Red gradient background (#FF6B6B to #FF8E8E)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['red_gradient_start']
    
    # Large countdown
    countdown_box = slide.shapes.add_textbox(Inches(2), Inches(1.5), Inches(9.33), Inches(2))
    countdown_frame = countdown_box.text_frame
    countdown_frame.clear()
    countdown_p = countdown_frame.paragraphs[0]
    countdown_p.text = "3"
    countdown_p.font.size = Pt(120)
    countdown_p.font.bold = True
    countdown_p.font.name = 'Roboto'
    countdown_p.font.color.rgb = colors['white']
    countdown_p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11.33), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.clear()
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.text = "That's All You Have"
    subtitle_p.font.size = Pt(36)
    subtitle_p.font.name = 'Roboto'
    subtitle_p.font.color.rgb = colors['white']
    subtitle_p.alignment = PP_ALIGN.CENTER
    
    # Statistics
    stats_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11.33), Inches(1.5))
    stats_frame = stats_box.text_frame
    stats_frame.clear()
    
    stats = [
        "53% of users abandon if load > 3 seconds",
        "4.47 MB initial bundle",
        "Performance Score: 45/100"
    ]
    
    for i, stat in enumerate(stats):
        p = stats_frame.paragraphs[i] if i < len(stats_frame.paragraphs) else stats_frame.add_paragraph()
        p.text = stat
        p.font.size = Pt(20)
        p.font.name = 'Roboto'
        p.font.color.rgb = colors['light_gray']
        p.alignment = PP_ALIGN.CENTER
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = "Don't shame the current state. Frame as 'opportunity for improvement'. Use hand gestures to emphasize the 3-second countdown."

def create_lighthouse_metrics_slide(prs, colors):
    """Create the Lighthouse metrics explanation slide with blue gradient background"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Blue gradient background (#4285F4 to #87CEEB)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['blue_gradient_start']
    
    # Title - Roboto Bold, 36px, White
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_p = title_frame.paragraphs[0]
    title_p.text = "What Lighthouse Really Measures"
    title_p.font.size = Pt(36)
    title_p.font.bold = True
    title_p.font.name = 'Roboto'
    title_p.font.color.rgb = colors['white']
    title_p.alignment = PP_ALIGN.CENTER
    
    # Metrics explanation
    metrics = [
        ("FCP", "First Contentful Paint", "When users see something", colors['performance_green']),
        ("LCP", "Largest Contentful Paint", "When they see what matters", colors['lighthouse_blue']),
        ("TTI", "Time to Interactive", "When they can do something", colors['orange_gradient_start']),
        ("CLS", "Cumulative Layout Shift", "Whether the page is stable", RGBColor(255, 100, 100))
    ]
    
    y_start = 2
    for i, (abbr, full, desc, color) in enumerate(metrics):
        # Metric abbreviation
        abbr_box = slide.shapes.add_textbox(Inches(1), Inches(y_start + i * 1.2), Inches(1.5), Inches(0.8))
        abbr_frame = abbr_box.text_frame
        abbr_frame.clear()
        abbr_p = abbr_frame.paragraphs[0]
        abbr_p.text = abbr
        abbr_p.font.size = Pt(24)
        abbr_p.font.bold = True
        abbr_p.font.color.rgb = color
        abbr_p.alignment = PP_ALIGN.CENTER
        
        # Full name
        name_box = slide.shapes.add_textbox(Inches(3), Inches(y_start + i * 1.2), Inches(4), Inches(0.4))
        name_frame = name_box.text_frame
        name_frame.clear()
        name_p = name_frame.paragraphs[0]
        name_p.text = full
        name_p.font.size = Pt(18)
        name_p.font.bold = True
        name_p.font.color.rgb = colors['white']
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(3), Inches(y_start + i * 1.2 + 0.4), Inches(8), Inches(0.6))
        desc_frame = desc_box.text_frame
        desc_frame.clear()
        desc_p = desc_frame.paragraphs[0]
        desc_p.text = desc
        desc_p.font.size = Pt(16)
        desc_p.font.color.rgb = colors['light_gray']
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = "Use hand gestures to show timeline progression. Emphasize that Lighthouse measures user experience, not just speed."

def create_approach_slide(prs, colors):
    """Create the optimization approach slide with green gradient background"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Green gradient background (#34A853 to #90EE90)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['green_gradient_start']
    
    # Title - Roboto Bold, 36px, White
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_p = title_frame.paragraphs[0]
    title_p.text = "Measure → Analyze → Optimize"
    title_p.font.size = Pt(36)
    title_p.font.bold = True
    title_p.font.name = 'Roboto'
    title_p.font.color.rgb = colors['white']
    title_p.alignment = PP_ALIGN.CENTER
    
    # Process steps
    steps = [
        ("1", "Measure", "Lighthouse CI baseline", colors['lighthouse_blue']),
        ("2", "Analyze", "Bundle analysis & bottleneck identification", colors['orange_gradient_start']),
        ("3", "Optimize", "Targeted interventions", colors['performance_green'])
    ]
    
    for i, (num, title, desc, color) in enumerate(steps):
        x_pos = 1 + i * 3.8
        
        # Step number
        num_box = slide.shapes.add_textbox(Inches(x_pos), Inches(2.5), Inches(1), Inches(1))
        num_frame = num_box.text_frame
        num_frame.clear()
        num_p = num_frame.paragraphs[0]
        num_p.text = num
        num_p.font.size = Pt(48)
        num_p.font.bold = True
        num_p.font.color.rgb = color
        num_p.alignment = PP_ALIGN.CENTER
        
        # Step title
        title_box = slide.shapes.add_textbox(Inches(x_pos), Inches(3.8), Inches(3), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.clear()
        title_p = title_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(20)
        title_p.font.bold = True
        title_p.font.color.rgb = colors['white']
        title_p.alignment = PP_ALIGN.CENTER
        
        # Step description
        desc_box = slide.shapes.add_textbox(Inches(x_pos), Inches(4.6), Inches(3), Inches(1))
        desc_frame = desc_box.text_frame
        desc_frame.clear()
        desc_p = desc_frame.paragraphs[0]
        desc_p.text = desc
        desc_p.font.size = Pt(14)
        desc_p.font.color.rgb = colors['light_gray']
        desc_p.alignment = PP_ALIGN.CENTER
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = "'No optimization without measurement.' 'We use Lighthouse CI to make performance a first-class citizen.'"

def create_lazy_loading_slide(prs, colors):
    """Create the lazy loading deep dive slide with orange gradient background"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Orange gradient background (#FBBC04 to #FFD700)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['orange_gradient_start']
    
    # Title - Roboto Bold, 36px, White
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_p = title_frame.paragraphs[0]
    title_p.text = "Lazy Loading: Load Only What You Need"
    title_p.font.size = Pt(36)
    title_p.font.bold = True
    title_p.font.name = 'Roboto'
    title_p.font.color.rgb = colors['white']
    title_p.alignment = PP_ALIGN.CENTER
    
    # Before/After code
    before_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(6), Inches(2.5))
    before_frame = before_box.text_frame
    before_frame.clear()
    before_p = before_frame.paragraphs[0]
    before_p.text = "Before: Eager Loading\n\nimport { CommitteeFormBuilderComponent }\nfrom './committee-form-builder';"
    before_p.font.size = Pt(14)
    before_p.font.color.rgb = colors['light_gray']
    before_p.font.name = 'Courier New'
    
    after_box = slide.shapes.add_textbox(Inches(6.5), Inches(2), Inches(6), Inches(2.5))
    after_frame = after_box.text_frame
    after_frame.clear()
    after_p = after_frame.paragraphs[0]
    after_p.text = "After: Route-level Lazy Loading\n\n{\n  path: 'CommitteeFormBuilder/:id',\n  loadChildren: () => import('./features/committee-builder.module')\n    .then(m => m.CommitteeBuilderModule)\n}"
    after_p.font.size = Pt(14)
    after_p.font.color.rgb = colors['performance_green']
    after_p.font.name = 'Courier New'
    
    # Impact
    impact_box = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(9.33), Inches(1.5))
    impact_frame = impact_box.text_frame
    impact_frame.clear()
    impact_p = impact_frame.paragraphs[0]
    impact_p.text = "Impact: Initial bundle reduced by ~792KB • TTI improved by ~2-3 seconds"
    impact_p.font.size = Pt(20)
    impact_p.font.bold = True
    impact_p.font.color.rgb = colors['orange_gradient_start']
    impact_p.alignment = PP_ALIGN.CENTER
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = "'We've already started this journey—now we finish it.' Show the dramatic bundle size reduction."

def create_onpush_slide(prs, colors):
    """Create the OnPush change detection slide with purple gradient background"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Purple gradient background (#9C27B0 to #E1BEE7)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['purple_gradient_start']
    
    # Title - Roboto Bold, 36px, White
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_p = title_frame.paragraphs[0]
    title_p.text = "OnPush: Smarter Change Detection"
    title_p.font.size = Pt(36)
    title_p.font.bold = True
    title_p.font.name = 'Roboto'
    title_p.font.color.rgb = colors['white']
    title_p.alignment = PP_ALIGN.CENTER
    
    # Problem statement
    problem_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(0.8))
    problem_frame = problem_box.text_frame
    problem_frame.clear()
    problem_p = problem_frame.paragraphs[0]
    problem_p.text = "The Problem: Every HTTP request triggers change detection across the entire component tree"
    problem_p.font.size = Pt(18)
    problem_p.font.color.rgb = colors['light_gray']
    problem_p.alignment = PP_ALIGN.CENTER
    
    # Code example
    code_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9.33), Inches(1.5))
    code_frame = code_box.text_frame
    code_frame.clear()
    code_p = code_frame.paragraphs[0]
    code_p.text = "@Component({\n  selector: 'app-dashboard',\n  templateUrl: './dashboard.component.html',\n  changeDetection: ChangeDetectionStrategy.OnPush\n})"
    code_p.font.size = Pt(14)
    code_p.font.color.rgb = colors['lighthouse_blue']
    code_p.font.name = 'Courier New'
    code_p.alignment = PP_ALIGN.CENTER
    
    # Impact
    impact_box = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(9.33), Inches(1))
    impact_frame = impact_box.text_frame
    impact_frame.clear()
    impact_p = impact_frame.paragraphs[0]
    impact_p.text = "Impact: 40-60% reduction in change detection cycles • TTI improvement: ~1-2 seconds"
    impact_p.font.size = Pt(18)
    impact_p.font.bold = True
    impact_p.font.color.rgb = colors['performance_green']
    impact_p.alignment = PP_ALIGN.CENTER
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = "'Not every component needs this—start with leaf components.' Show the dramatic reduction in CPU usage."

def create_font_loading_slide(prs, colors):
    """Create the font loading strategy slide with teal gradient background"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Teal gradient background (#009688 to #80CBC4)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['teal_gradient_start']
    
    # Title - Roboto Bold, 36px, White
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_p = title_frame.paragraphs[0]
    title_p.text = "Async Font Loading: Don't Block the Critical Path"
    title_p.font.size = Pt(36)
    title_p.font.bold = True
    title_p.font.name = 'Roboto'
    title_p.font.color.rgb = colors['white']
    title_p.alignment = PP_ALIGN.CENTER
    
    # Problem
    problem_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(0.6))
    problem_frame = problem_box.text_frame
    problem_frame.clear()
    problem_p = problem_frame.paragraphs[0]
    problem_p.text = "Current Problem: These block rendering!"
    problem_p.font.size = Pt(18)
    problem_p.font.color.rgb = colors['light_gray']
    problem_p.alignment = PP_ALIGN.CENTER
    
    # Problem code
    problem_code_box = slide.shapes.add_textbox(Inches(2), Inches(2.8), Inches(9.33), Inches(0.8))
    problem_code_frame = problem_code_box.text_frame
    problem_code_frame.clear()
    problem_code_p = problem_code_frame.paragraphs[0]
    problem_code_p.text = '<link href="https://fonts.googleapis.com/css2?family=Roboto:wght@300;400;500&display=swap" rel="stylesheet">'
    problem_code_p.font.size = Pt(12)
    problem_code_p.font.color.rgb = RGBColor(255, 100, 100)
    problem_code_p.font.name = 'Courier New'
    problem_code_p.alignment = PP_ALIGN.CENTER
    
    # Solution
    solution_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.33), Inches(0.6))
    solution_frame = solution_box.text_frame
    solution_frame.clear()
    solution_p = solution_frame.paragraphs[0]
    solution_p.text = "Optimized Approach: Preconnect + Async Load"
    solution_p.font.size = Pt(18)
    solution_p.font.color.rgb = colors['performance_green']
    solution_p.alignment = PP_ALIGN.CENTER
    
    # Solution code
    solution_code_box = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(11.33), Inches(1.2))
    solution_code_frame = solution_code_box.text_frame
    solution_code_frame.clear()
    solution_code_p = solution_code_frame.paragraphs[0]
    solution_code_p.text = '<!-- Preconnect to speed up DNS/TCP/TLS -->\n<link rel="preconnect" href="https://fonts.googleapis.com">\n<link rel="preconnect" href="https://fonts.gstatic.com" crossorigin>\n\n<!-- Async load with font-display: swap -->\n<link href="https://fonts.googleapis.com/css2?family=Roboto:wght@300;400;500&display=swap"\n      rel="stylesheet" media="print" onload="this.media=\'all\'">'
    solution_code_p.font.size = Pt(10)
    solution_code_p.font.color.rgb = colors['accent_cyan']
    solution_code_p.font.name = 'Courier New'
    solution_code_p.alignment = PP_ALIGN.CENTER
    
    # Impact
    impact_box = slide.shapes.add_textbox(Inches(2), Inches(6.2), Inches(9.33), Inches(0.8))
    impact_frame = impact_box.text_frame
    impact_frame.clear()
    impact_p = impact_frame.paragraphs[0]
    impact_p.text = "Impact: FCP improvement of 400-800ms • Better perceived performance"
    impact_p.font.size = Pt(16)
    impact_p.font.bold = True
    impact_p.font.color.rgb = colors['orange_gradient_start']
    impact_p.alignment = PP_ALIGN.CENTER
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = "Show the dramatic difference between blocking and non-blocking font loading. Emphasize the FCP improvement."

def create_material_modules_slide(prs, colors):
    """Create the tree-shakeable Material modules slide with indigo gradient background"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Indigo gradient background (#3F51B5 to #9FA8DA)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['indigo_gradient_start']
    
    # Title - Roboto Bold, 36px, White
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_p = title_frame.paragraphs[0]
    title_p.text = "Refactor SharedModule: Import Only What You Use"
    title_p.font.size = Pt(36)
    title_p.font.bold = True
    title_p.font.name = 'Roboto'
    title_p.font.color.rgb = colors['white']
    title_p.alignment = PP_ALIGN.CENTER
    
    # Problem
    problem_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(0.6))
    problem_frame = problem_box.text_frame
    problem_frame.clear()
    problem_p = problem_frame.paragraphs[0]
    problem_p.text = "The Problem: SharedModule imports and exports 25+ Material modules"
    problem_p.font.size = Pt(18)
    problem_p.font.color.rgb = colors['light_gray']
    problem_p.alignment = PP_ALIGN.CENTER
    
    # Code example
    code_box = slide.shapes.add_textbox(Inches(2), Inches(2.8), Inches(9.33), Inches(1.5))
    code_frame = code_box.text_frame
    code_frame.clear()
    code_p = code_frame.paragraphs[0]
    code_p.text = "// shared.module.ts - Importing EVERYTHING\nimports: [\n  MatDialogModule, MatProgressSpinnerModule, MatIconModule,\n  MatCardModule, MatDatepickerModule, MatFormFieldModule,\n  // ... 20 more modules\n]"
    code_p.font.size = Pt(12)
    code_p.font.color.rgb = RGBColor(255, 100, 100)
    code_p.font.name = 'Courier New'
    code_p.alignment = PP_ALIGN.CENTER
    
    # Solution
    solution_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.33), Inches(0.6))
    solution_frame = solution_box.text_frame
    solution_frame.clear()
    solution_p = solution_frame.paragraphs[0]
    solution_p.text = "Optimized Approach: Feature-specific imports only"
    solution_p.font.size = Pt(18)
    solution_p.font.color.rgb = colors['performance_green']
    solution_p.alignment = PP_ALIGN.CENTER
    
    # Impact
    impact_box = slide.shapes.add_textbox(Inches(2), Inches(5.5), Inches(9.33), Inches(0.8))
    impact_frame = impact_box.text_frame
    impact_frame.clear()
    impact_p = impact_frame.paragraphs[0]
    impact_p.text = "Impact: ~200-300KB reduction in vendor bundle • Better tree-shaking"
    impact_p.font.size = Pt(16)
    impact_p.font.bold = True
    impact_p.font.color.rgb = colors['orange_gradient_start']
    impact_p.alignment = PP_ALIGN.CENTER
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = "Show the dramatic difference in bundle size. Emphasize that this is about importing only what you actually use."

def create_demo_slide(prs, colors):
    """Create the live demo slide with dark gradient background"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Dark gradient background (#2C3E50 to #34495E)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['dark_gradient_start']
    
    # Title - Roboto Bold, 36px, White
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_p = title_frame.paragraphs[0]
    title_p.text = "Watch the Magic Happen"
    title_p.font.size = Pt(36)
    title_p.font.bold = True
    title_p.font.name = 'Roboto'
    title_p.font.color.rgb = colors['white']
    title_p.alignment = PP_ALIGN.CENTER
    
    # Demo script
    demo_box = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(9.33), Inches(3))
    demo_frame = demo_box.text_frame
    demo_frame.clear()
    
    demo_steps = [
        "1. Here's the current dashboard—watch the metrics",
        "2. I'll apply OnPush to just 3 components",
        "3. Rebuild and re-run Lighthouse",
        "4. Notice the TTI dropped from 8.2s to 5.7s"
    ]
    
    for i, step in enumerate(demo_steps):
        p = demo_frame.paragraphs[i] if i < len(demo_frame.paragraphs) else demo_frame.add_paragraph()
        p.text = step
        p.font.size = Pt(18)
        p.font.color.rgb = colors['light_gray']
        p.alignment = PP_ALIGN.CENTER
    
    # Performance indicator
    perf_box = slide.shapes.add_textbox(Inches(3), Inches(5.5), Inches(7.33), Inches(1))
    perf_frame = perf_box.text_frame
    perf_frame.clear()
    perf_p = perf_frame.paragraphs[0]
    perf_p.text = "TTI: 8.2s → 5.7s"
    perf_p.font.size = Pt(24)
    perf_p.font.bold = True
    perf_p.font.color.rgb = colors['performance_green']
    perf_p.alignment = PP_ALIGN.CENTER
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = "Have a backup video in case of technical difficulties. Show the dramatic improvement in real-time."

def create_results_slide(prs, colors):
    """Create the before/after results slide with success gradient background"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Success gradient background (#27AE60 to #58D68D)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['success_gradient_start']
    
    # Title - Roboto Bold, 36px, White
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_p = title_frame.paragraphs[0]
    title_p.text = "The Numbers Don't Lie"
    title_p.font.size = Pt(36)
    title_p.font.bold = True
    title_p.font.name = 'Roboto'
    title_p.font.color.rgb = colors['white']
    title_p.alignment = PP_ALIGN.CENTER
    
    # Results table
    results = [
        ("Performance Score", "45", "78", "+73%"),
        ("FCP", "2.8s", "1.2s", "-57%"),
        ("LCP", "6.4s", "2.3s", "-64%"),
        ("TTI", "8.2s", "3.1s", "-62%"),
        ("Bundle Size", "4.47 MB", "2.1 MB", "-53%")
    ]
    
    # Table headers
    headers = ["Metric", "Before", "After", "Improvement"]
    for i, header in enumerate(headers):
        header_box = slide.shapes.add_textbox(Inches(1 + i * 2.8), Inches(2), Inches(2.8), Inches(0.6))
        header_frame = header_box.text_frame
        header_frame.clear()
        header_p = header_frame.paragraphs[0]
        header_p.text = header
        header_p.font.size = Pt(16)
        header_p.font.bold = True
        header_p.font.color.rgb = colors['orange_gradient_start']
        header_p.alignment = PP_ALIGN.CENTER
    
    # Table rows
    for row_idx, (metric, before, after, improvement) in enumerate(results):
        y_pos = 2.8 + row_idx * 0.6
        
        # Metric name
        metric_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(2.8), Inches(0.6))
        metric_frame = metric_box.text_frame
        metric_frame.clear()
        metric_p = metric_frame.paragraphs[0]
        metric_p.text = metric
        metric_p.font.size = Pt(14)
        metric_p.font.color.rgb = colors['white']
        metric_p.alignment = PP_ALIGN.CENTER
        
        # Before value
        before_box = slide.shapes.add_textbox(Inches(3.8), Inches(y_pos), Inches(2.8), Inches(0.6))
        before_frame = before_box.text_frame
        before_frame.clear()
        before_p = before_frame.paragraphs[0]
        before_p.text = before
        before_p.font.size = Pt(14)
        before_p.font.color.rgb = RGBColor(255, 100, 100)
        before_p.alignment = PP_ALIGN.CENTER
        
        # After value
        after_box = slide.shapes.add_textbox(Inches(6.6), Inches(y_pos), Inches(2.8), Inches(0.6))
        after_frame = after_box.text_frame
        after_frame.clear()
        after_p = after_frame.paragraphs[0]
        after_p.text = after
        after_p.font.size = Pt(14)
        after_p.font.color.rgb = colors['performance_green']
        after_p.alignment = PP_ALIGN.CENTER
        
        # Improvement
        improvement_box = slide.shapes.add_textbox(Inches(9.4), Inches(y_pos), Inches(2.8), Inches(0.6))
        improvement_frame = improvement_box.text_frame
        improvement_frame.clear()
        improvement_p = improvement_frame.paragraphs[0]
        improvement_p.text = improvement
        improvement_p.font.size = Pt(14)
        improvement_p.font.bold = True
        improvement_p.font.color.rgb = colors['accent_cyan']
        improvement_p.alignment = PP_ALIGN.CENTER
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = "'These aren't theoretical—these are real improvements.' Emphasize the dramatic improvements."

def create_roadmap_slide(prs, colors):
    """Create the future roadmap slide with future gradient background"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Future gradient background (#8E44AD to #BB8FCE)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['future_gradient_start']
    
    # Title - Roboto Bold, 36px, White
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_p = title_frame.paragraphs[0]
    title_p.text = "The Journey Continues"
    title_p.font.size = Pt(36)
    title_p.font.bold = True
    title_p.font.name = 'Roboto'
    title_p.font.color.rgb = colors['white']
    title_p.alignment = PP_ALIGN.CENTER
    
    # High Impact / Low Effort
    high_low_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(5.5), Inches(2))
    high_low_frame = high_low_box.text_frame
    high_low_frame.clear()
    high_low_p = high_low_frame.paragraphs[0]
    high_low_p.text = "High Impact / Low Effort (Do First)"
    high_low_p.font.size = Pt(18)
    high_low_p.font.bold = True
    high_low_p.font.color.rgb = colors['performance_green']
    
    high_low_items = [
        "• Image lazy loading with loading='lazy'",
        "• Remove unused polyfills",
        "• Enable Brotli compression"
    ]
    
    for i, item in enumerate(high_low_items):
        p = high_low_frame.paragraphs[i + 1] if i + 1 < len(high_low_frame.paragraphs) else high_low_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = colors['light_gray']
    
    # High Impact / High Effort
    high_high_box = slide.shapes.add_textbox(Inches(6.5), Inches(2), Inches(5.5), Inches(2))
    high_high_frame = high_high_box.text_frame
    high_high_frame.clear()
    high_high_p = high_high_frame.paragraphs[0]
    high_high_p.text = "High Impact / High Effort (Plan Carefully)"
    high_high_p.font.size = Pt(18)
    high_high_p.font.bold = True
    high_high_p.font.color.rgb = colors['accent_orange']
    
    high_high_items = [
        "• Migrate to standalone components (Angular 18+)",
        "• Implement virtual scrolling for large lists",
        "• Progressive Web App (PWA) features"
    ]
    
    for i, item in enumerate(high_high_items):
        p = high_high_frame.paragraphs[i + 1] if i + 1 < len(high_high_frame.paragraphs) else high_high_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = colors['light_gray']
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = "Show the strategic approach to optimization. Emphasize that this is just the beginning."

def create_cta_slide(prs, colors):
    """Create the call to action slide with action gradient background"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Action gradient background (#E74C3C to #F1948A)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['action_gradient_start']
    
    # Title - Roboto Bold, 36px, White
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_p = title_frame.paragraphs[0]
    title_p.text = "Performance is a Feature, Not a Nice-to-Have"
    title_p.font.size = Pt(36)
    title_p.font.bold = True
    title_p.font.name = 'Roboto'
    title_p.font.color.rgb = colors['white']
    title_p.alignment = PP_ALIGN.CENTER
    
    # Actionable takeaways
    takeaways_box = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(9.33), Inches(2))
    takeaways_frame = takeaways_box.text_frame
    takeaways_frame.clear()
    
    takeaways = [
        "1. Add Lighthouse CI to your pipeline today",
        "2. Set budget limits in angular.json",
        "3. Make one optimization per sprint"
    ]
    
    for i, takeaway in enumerate(takeaways):
        p = takeaways_frame.paragraphs[i] if i < len(takeaways_frame.paragraphs) else takeaways_frame.add_paragraph()
        p.text = takeaway
        p.font.size = Pt(20)
        p.font.color.rgb = colors['light_gray']
        p.alignment = PP_ALIGN.CENTER
    
    # Closing line
    closing_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.33), Inches(1.5))
    closing_frame = closing_box.text_frame
    closing_frame.clear()
    closing_p = closing_frame.paragraphs[0]
    closing_p.text = "Remember: Every millisecond is a vote of confidence in your users."
    closing_p.font.size = Pt(24)
    closing_p.font.bold = True
    closing_p.font.color.rgb = colors['accent_orange']
    closing_p.alignment = PP_ALIGN.CENTER
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = "End with energy and conviction. Make it clear this is achievable. Encourage immediate action."

def create_qa_slide(prs, colors):
    """Create the Q&A slide with discussion gradient background"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Discussion gradient background (#3498DB to #85C1E9)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['discussion_gradient_start']
    
    # Title - Roboto Bold, 36px, White
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_p = title_frame.paragraphs[0]
    title_p.text = "Questions & Discussion"
    title_p.font.size = Pt(36)
    title_p.font.bold = True
    title_p.font.name = 'Roboto'
    title_p.font.color.rgb = colors['white']
    title_p.alignment = PP_ALIGN.CENTER
    
    # Resources
    resources_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9.33), Inches(2))
    resources_frame = resources_box.text_frame
    resources_frame.clear()
    resources_p = resources_frame.paragraphs[0]
    resources_p.text = "Key Resources:"
    resources_p.font.size = Pt(20)
    resources_p.font.bold = True
    resources_p.font.color.rgb = colors['accent_cyan']
    resources_p.alignment = PP_ALIGN.CENTER
    
    resources = [
        "• Angular Performance Guide",
        "• Lighthouse Documentation",
        "• Web.dev Performance",
        "• Core Web Vitals"
    ]
    
    for i, resource in enumerate(resources):
        p = resources_frame.paragraphs[i + 1] if i + 1 < len(resources_frame.paragraphs) else resources_frame.add_paragraph()
        p.text = resource
        p.font.size = Pt(16)
        p.font.color.rgb = colors['light_gray']
        p.alignment = PP_ALIGN.CENTER
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = "Be prepared for technical questions. Have backup slides with more details. Encourage discussion and feedback."

def create_thank_you_slide(prs, colors):
    """Create the thank you slide with gratitude gradient background"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Gratitude gradient background (#F39C12 to #F7DC6F)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['gratitude_gradient_start']
    
    # Title - Roboto Bold, 36px, White
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_p = title_frame.paragraphs[0]
    title_p.text = "Thank You for Your Time"
    title_p.font.size = Pt(36)
    title_p.font.bold = True
    title_p.font.name = 'Roboto'
    title_p.font.color.rgb = colors['white']
    title_p.alignment = PP_ALIGN.CENTER
    
    # Key message
    message_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.33), Inches(1.5))
    message_frame = message_box.text_frame
    message_frame.clear()
    message_p = message_frame.paragraphs[0]
    message_p.text = "Let's make the web faster, one optimization at a time."
    message_p.font.size = Pt(24)
    message_p.font.color.rgb = colors['accent_orange']
    message_p.alignment = PP_ALIGN.CENTER
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = "End on a positive note. Encourage follow-up conversations. Thank the audience for their attention."

def main():
    """Main function to create the presentation"""
    print("Creating Angular Optimization TED Talk presentation...")
    
    # Create the presentation
    prs = create_angular_presentation()
    
    # Save the presentation
    output_path = "/workspace/Angular_Optimization_TED_Talk.pptx"
    prs.save(output_path)
    
    print(f"Presentation created successfully: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
    
    return output_path

if __name__ == "__main__":
    main()